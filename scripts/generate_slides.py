#!/usr/bin/env python3
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_FAMILY = "Helvetica Neue Light"
COLOR_BG = RGBColor(255, 255, 255)
COLOR_TEXT = RGBColor(0, 0, 0)
COLOR_ACCENT = RGBColor(0, 0, 0)


def load_slides(json_path: Path) -> list[dict]:
    with json_path.open("r", encoding="utf-8") as f:
        return json.load(f)


def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_border(slide):
    border = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.08),
        Inches(0.1),
        Inches(13.2),
        Inches(7.3),
    )
    border.fill.background()
    border.line.color.rgb = COLOR_ACCENT
    border.line.width = Pt(0.75)
    border.shadow.inherit = False


def add_title(slide, text: str, *, color=COLOR_TEXT, top=0.3, left=0.5, width=9.0, height=1.25, size=34):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(0)
    p.line_spacing = 1.0

    run = p.add_run()
    run.text = text
    run.font.name = FONT_FAMILY
    run.font.size = Pt(size)
    run.font.bold = False
    run.font.italic = False
    run.font.color.rgb = color
    return box


def add_body(slide, text: str, *, left=0.5, top=1.75, width=9.0, height=4.95, font_size=16, color=COLOR_TEXT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    lines = text.split("\n")
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = line
        p.font.name = FONT_FAMILY
        p.font.size = Pt(font_size)
        p.font.bold = False
        p.font.italic = False
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.space_after = Pt(0)
        p.line_spacing = 1.0
    return box


def add_code_block(slide, code: str):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.75), Inches(9.0), Inches(4.95))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    lines = code.split("\n")
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = "Menlo"
        run.font.size = Pt(16)
        run.font.bold = False
        run.font.italic = False
        run.font.color.rgb = COLOR_TEXT
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)
        p.line_spacing = 1.0


def build_title_slide(prs: Presentation, slide_data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLOR_BG)
    add_border(slide)

    box = slide.shapes.add_textbox(Inches(1.0), Inches(3.3), Inches(11.4), Inches(2.5))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(0)
    p.line_spacing = 1.0

    run = p.add_run()
    run.text = slide_data.get("content", "") or slide_data.get("title", "")
    run.font.name = FONT_FAMILY
    run.font.size = Pt(32)
    run.font.bold = False
    run.font.italic = False
    run.font.color.rgb = COLOR_TEXT


def build_content_slide(prs: Presentation, slide_data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLOR_BG)
    add_border(slide)
    add_title(slide, slide_data.get("title", ""))
    add_body(slide, slide_data.get("content", ""))


def build_code_slide(prs: Presentation, slide_data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, COLOR_BG)
    add_border(slide)
    add_title(slide, slide_data.get("title", ""))
    add_code_block(slide, slide_data.get("content", ""))


def build_presentation(slides_data: list[dict], output_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for slide_data in slides_data:
        slide_type = slide_data.get("type", "content")
        if slide_type == "title":
            build_title_slide(prs, slide_data)
        elif slide_type == "code":
            build_code_slide(prs, slide_data)
        else:
            build_content_slide(prs, slide_data)

        notes_text = slide_data.get("speaker_notes", "")
        if notes_text:
            notes_frame = prs.slides[-1].notes_slide.notes_text_frame
            notes_frame.text = notes_text

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def main():
    if len(sys.argv) != 3:
        print("Usage: python scripts/generate_slides.py <slides_content.json> <output.pptx>")
        sys.exit(1)

    input_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2])

    slides_data = load_slides(input_path)
    build_presentation(slides_data, output_path)
    print(f"Wrote {output_path}")


if __name__ == "__main__":
    main()
